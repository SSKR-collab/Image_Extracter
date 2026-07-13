import os
import json
from PIL import Image
from image_extractor.base_analyzer import BaseAnalyzer

# Check for Tesseract and EasyOCR dependencies
try:
    import pytesseract
    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False

try:
    import easyocr
    HAS_EASYOCR = True
except ImportError:
    HAS_EASYOCR = False

try:
    import cv2
    import numpy as np
    HAS_OPENCV = True
except ImportError:
    HAS_OPENCV = False


class OcrEngine(BaseAnalyzer):
    """
    Coordinates Optical Character Recognition. Supports EasyOCR, Tesseract,
    and falls back to loading sidecar (.ocr/.txt/.json) files for offline/test environments.
    Implements advanced local Tesseract preprocessing (tiling, multi-scale, multi-PSM, CLAHE, sharpening)
    to maximize text extraction accuracy on complex layouts without heavy models.
    """
    VERSION = "1.3.0"

    def analyze(self, file_path: str, img, context: dict) -> dict:
        results = {
            "facts": {
                "raw_text": "",
                "words": []
            },
            "indicators": [],
            "assessments": {},
            "errors": []
        }

        # Format-based extraction routing
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return self._analyze_pdf(file_path, results)
        elif ext == ".docx":
            return self._analyze_docx(file_path, results)
        elif ext == ".pptx":
            return self._analyze_pptx(file_path, results)
        elif ext in (".xlsx", ".xls"):
            return self._analyze_xlsx(file_path, results)
        elif ext in (".txt", ".md", ".json", ".csv", ".xml", ".html", ".yaml", ".yml"):
            return self._analyze_text_file(file_path, results)

        # Check sidecar fallback first
        base_path, _ = os.path.splitext(file_path)
        sidecar_ocr = base_path + ".ocr"
        sidecar_txt = base_path + ".txt"
        sidecar_json = base_path + ".json"
        
        selected_sidecar = None
        if os.path.exists(sidecar_ocr):
            selected_sidecar = sidecar_ocr
        elif os.path.exists(sidecar_txt):
            selected_sidecar = sidecar_txt
        elif os.path.exists(sidecar_json):
            selected_sidecar = sidecar_json

        if selected_sidecar:
            try:
                with open(selected_sidecar, "r", encoding="utf-8") as sf:
                    content = sf.read().strip()
                
                # Check if sidecar is structured JSON
                try:
                    json_data = json.loads(content)
                    if isinstance(json_data, dict):
                        raw_text = ""
                        words_list = []
                        
                        # Support multiple key variations
                        if "raw_text" in json_data:
                            raw_text = json_data["raw_text"]
                        elif "ocr_data" in json_data:
                            if isinstance(json_data["ocr_data"], dict):
                                raw_text = json_data["ocr_data"].get("raw_text", "")
                                words_list = json_data["ocr_data"].get("words", [])
                            else:
                                raw_text = str(json_data["ocr_data"])
                        elif "facts" in json_data and "raw_text" in json_data["facts"]:
                            raw_text = json_data["facts"]["raw_text"]
                            words_list = json_data["facts"].get("words", [])
                            
                        if "words" in json_data:
                            words_list = json_data["words"]
                            
                        if raw_text or words_list:
                            results["facts"]["raw_text"] = raw_text
                            results["facts"]["words"] = words_list if words_list else self._mock_words_from_text(raw_text)
                        else:
                            raise ValueError()
                    else:
                        raise ValueError()
                except Exception:
                    # Plain text sidecar
                    results["facts"]["raw_text"] = content
                    results["facts"]["words"] = self._mock_words_from_text(content)

                results["indicators"].append({
                    "type": "ocr_sidecar_loaded",
                    "description": f"Loaded OCR sidecar file: {os.path.basename(selected_sidecar)}",
                    "severity": "low"
                })
                return results
            except Exception as e:
                results["errors"].append({
                    "plugin": self.get_name(),
                    "severity": "warning",
                    "message": f"Failed to load OCR sidecar: {str(e)}"
                })

        # Try EasyOCR
        if HAS_EASYOCR:
            try:
                # Convert PIL Image to byte buffer
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format=img.format or 'PNG')
                img_bytes = img_byte_arr.getvalue()
                
                reader = easyocr.Reader(['en'], gpu=False)
                ocr_results = reader.readtext(img_bytes)
                
                raw_lines = []
                words_list = []
                for bbox, text, conf in ocr_results:
                    raw_lines.append(text)
                    xs = [pt[0] for pt in bbox]
                    ys = [pt[1] for pt in bbox]
                    x_min, y_min = min(xs), min(ys)
                    w, h = max(xs) - x_min, max(ys) - y_min
                    
                    words_list.append({
                        "text": text,
                        "bbox": [x_min, y_min, w, h],
                        "confidence": round(float(conf), 3)
                    })
                
                results["facts"]["raw_text"] = "\n".join(raw_lines)
                results["facts"]["words"] = words_list
                return results
            except Exception as e:
                results["errors"].append({
                    "plugin": self.get_name(),
                    "severity": "warning",
                    "message": f"EasyOCR execution failed: {str(e)}"
                })

        # Try Tesseract
        if HAS_TESSERACT:
            try:
                tess_path = self.config.get("tesseract_path")
                if tess_path:
                    pytesseract.pytesseract.tesseract_cmd = tess_path
                elif os.name == "nt":
                    default_paths = [
                        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                        os.path.expandvars(r"%LOCALAPPDATA%\Programs\Tesseract-OCR\tesseract.exe")
                    ]
                    for path in default_paths:
                        if os.path.exists(path):
                            pytesseract.pytesseract.tesseract_cmd = path
                            break

                # Tesseract OCR execution
                # Check if we can run advanced local preprocessing
                if HAS_OPENCV:
                    words_list = self._run_tesseract_with_enhancements(img, pytesseract.pytesseract.tesseract_cmd)
                else:
                    data_str = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
                    words_list = self._parse_tesseract_dict(data_str)

                # Sort words by line_num and then by x coordinate
                words_list.sort(key=lambda w: (w.get("line_num", 0), w["bbox"][0]))
                
                raw_lines = []
                current_line = []
                last_line_num = -1
                for w in words_list:
                    text = w["text"]
                    line_num = w.get("line_num", 0)
                    if line_num != last_line_num:
                        if current_line:
                            raw_lines.append(" ".join(current_line))
                        current_line = [text]
                        last_line_num = line_num
                    else:
                        current_line.append(text)
                if current_line:
                    raw_lines.append(" ".join(current_line))
                
                results["facts"]["raw_text"] = "\n".join(raw_lines)
                
                # Clean line_num helper keys
                for w in words_list:
                    if "line_num" in w:
                        del w["line_num"]
                
                results["facts"]["words"] = words_list
                return results
            except Exception as e:
                results["errors"].append({
                    "plugin": self.get_name(),
                    "severity": "warning",
                    "message": f"Tesseract execution failed: {str(e)}"
                })

        # If no OCR is available
        results["errors"].append({
            "plugin": self.get_name(),
            "severity": "warning",
            "message": "No local OCR engine (EasyOCR or Tesseract) is installed, and no sidecar OCR file was found."
        })
        return results

    def _run_tesseract_with_enhancements(self, img, tess_cmd) -> list:
        """
        Runs Tesseract OCR using multiple scales, PSM config values, and
        OpenCV image preprocessing filters to maximize extraction accuracy.
        """
        # Convert Pillow image to OpenCV BGR format
        open_cv_image = np.array(img.convert("RGB"))
        gray = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2GRAY)
        
        # Pass 1: Run raw image first
        data_str = pytesseract.image_to_data(img, output_type=pytesseract.Output.DICT)
        words = self._parse_tesseract_dict(data_str)
        
        # Pass 2: Enhanced Grayscale (CLAHE + Sharpening)
        try:
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
            contrast = clahe.apply(gray)
            kernel = np.array([[0, -1, 0], [-1, 5, -1], [0, -1, 0]])
            sharpened = cv2.filter2D(contrast, -1, kernel)
            data_str_2 = pytesseract.image_to_data(Image.fromarray(sharpened), output_type=pytesseract.Output.DICT)
            words_2 = self._parse_tesseract_dict(data_str_2)
            
            for w2 in words_2:
                if not any(self._is_overlapping(w2["bbox"], w["bbox"]) for w in words):
                    words.append(w2)
        except Exception:
            pass
            
        # Pass 3: Adaptive Thresholding (brings out low-contrast background text)
        try:
            thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 15, 8)
            data_str_3 = pytesseract.image_to_data(Image.fromarray(thresh), output_type=pytesseract.Output.DICT)
            words_3 = self._parse_tesseract_dict(data_str_3)
            
            for w3 in words_3:
                if not any(self._is_overlapping(w3["bbox"], w["bbox"]) for w in words):
                    words.append(w3)
        except Exception:
            pass
            
        # If text is extremely short, run advanced tile-based preprocessing passes
        if len(words) < 15:
            # Attempt grid-based overlapping tile scans (captures small, curved, and stylized text segments)
            tile_words = self._run_tile_based_ocr(gray)
            for tw in tile_words:
                if not any(self._is_overlapping(tw["bbox"], w["bbox"]) for w in words):
                    words.append(tw)
                    
            try:
                # Rotation pass: test 90-degrees clockwise rotation for vertical/skewed text banners
                h, w = gray.shape
                rotated_90 = cv2.rotate(gray, cv2.ROTATE_90_CLOCKWISE)
                data_str_90 = pytesseract.image_to_data(Image.fromarray(rotated_90), output_type=pytesseract.Output.DICT)
                words_90 = self._parse_tesseract_dict(data_str_90)
                
                for w90 in words_90:
                    x, y, width, height = w90["bbox"]
                    # Map coordinates back: x_new = y, y_new = h - x - width
                    orig_x = y
                    orig_y = h - x - width
                    w90["bbox"] = [orig_x, orig_y, height, width]
                    
                    if not any(self._is_overlapping(w90["bbox"], w["bbox"]) for w in words):
                        words.append(w90)
            except Exception:
                pass
                
        # Split words joined by colons (e.g. "in:schools" -> "in", "schools")
        cleaned_words = []
        for w in words:
            txt = w["text"]
            if ":" in txt and not txt.startswith("http") and not txt.replace(":", "").isdigit():
                parts = txt.split(":")
                parts = [p for p in parts if p.strip()]
                if len(parts) > 1:
                    x, y, width, height = w["bbox"]
                    part_w = width // len(parts)
                    for idx, part in enumerate(parts):
                        cleaned_words.append({
                            "text": part,
                            "bbox": [x + idx * part_w, y, part_w, height],
                            "confidence": w["confidence"],
                            "line_num": w.get("line_num", 0)
                        })
                    continue
            cleaned_words.append(w)
            
        return cleaned_words

    def _run_tile_based_ocr(self, gray_img) -> list:
        """
        Divides the image into overlapping tiles and runs upscaled + enhanced
        multi-PSM Tesseract checks on each crop region.
        """
        h, w = gray_img.shape
        tile_size = 512
        step_size = 384
        
        words = []
        
        # Grid bounds calculation
        y_steps = list(range(0, max(1, h - tile_size), step_size))
        if not y_steps or y_steps[-1] + tile_size < h:
            y_steps.append(max(0, h - tile_size))
            
        x_steps = list(range(0, max(1, w - tile_size), step_size))
        if not x_steps or x_steps[-1] + tile_size < w:
            x_steps.append(max(0, w - tile_size))
            
        # Run sliding window
        for y in y_steps:
            for x in x_steps:
                tile = gray_img[y:y+tile_size, x:x+tile_size]
                th, tw = tile.shape
                if th == 0 or tw == 0:
                    continue
                    
                # Preprocess tile: Upscale 3x, CLAHE contrast stretch, filter sharpening
                resized_tile = cv2.resize(tile, (tw * 3, th * 3), interpolation=cv2.INTER_CUBIC)
                
                clahe = cv2.createCLAHE(clipLimit=2.5, tileGridSize=(8, 8))
                contrast_tile = clahe.apply(resized_tile)
                
                kernel = np.array([[0, -1, 0], [-1, 5, -1], [0, -1, 0]])
                sharpened_tile = cv2.filter2D(contrast_tile, -1, kernel)
                
                pil_tile = Image.fromarray(sharpened_tile)
                
                # Check multiple PSMs (default, sparse text, uniform block)
                for psm in ["3", "11", "6"]:
                    custom_config = f"--oem 3 --psm {psm}"
                    try:
                        data_str_tile = pytesseract.image_to_data(
                            pil_tile, config=custom_config, output_type=pytesseract.Output.DICT
                        )
                        tile_words = self._parse_tesseract_dict(data_str_tile)
                        
                        # Downscale bounding boxes back and translate to image space
                        for tw_word in tile_words:
                            tw_word["bbox"] = [
                                x + (tw_word["bbox"][0] // 3),
                                y + (tw_word["bbox"][1] // 3),
                                tw_word["bbox"][2] // 3,
                                tw_word["bbox"][3] // 3
                            ]
                            
                        # Add word if it doesn't overlap an existing detection
                        for tw_word in tile_words:
                            if not any(self._is_overlapping(tw_word["bbox"], mw["bbox"]) for mw in words):
                                words.append(tw_word)
                    except Exception:
                        pass
        return words

    def _parse_tesseract_dict(self, data_str) -> list:
        words = []
        n_boxes = len(data_str["text"])
        for i in range(n_boxes):
            text = data_str["text"][i].strip()
            if not text:
                continue
            
            conf = float(data_str["conf"][i]) / 100.0 if "conf" in data_str else 0.8
            if conf < 0:
                conf = 0.8
                
            x = data_str["left"][i]
            y = data_str["top"][i]
            w = data_str["width"][i]
            h = data_str["height"][i]
            
            words.append({
                "text": text,
                "bbox": [x, y, w, h],
                "confidence": round(conf, 3),
                "line_num": data_str.get("line_num", [0]*n_boxes)[i]
            })
        return words

    def _is_overlapping(self, boxA, boxB) -> bool:
        # boxA/boxB format: [x, y, w, h]
        xA = max(boxA[0], boxB[0])
        yA = max(boxA[1], boxB[1])
        xB = min(boxA[0] + boxA[2], boxB[0] + boxB[2])
        yB = min(boxA[1] + boxA[3], boxB[1] + boxB[3])
        
        interArea = max(0, xB - xA) * max(0, yB - yA)
        if interArea > 0:
            areaA = boxA[2] * boxA[3]
            areaB = boxB[2] * boxB[3]
            # Overlap ratio relative to smaller box
            overlap = interArea / float(min(areaA, areaB))
            return overlap > 0.4
        return False

    def _mock_words_from_text(self, text: str) -> list:
        """
        Generates simulated words with structured coordinates so that downstream
        layout, entities, and relationship analyzers can run on raw text inputs.
        """
        words = []
        lines = text.split("\n")
        
        y_offset = 20
        for line in lines:
            line_words = line.strip().split()
            if not line_words:
                y_offset += 30
                continue
                
            x_offset = 20
            for lw in line_words:
                w_len = len(lw) * 10
                words.append({
                    "text": lw,
                    "bbox": [x_offset, y_offset, w_len, 20],
                    "confidence": 0.95
                })
                x_offset += w_len + 15
            y_offset += 35
        return words

    def _rebuild_text_from_words(self, words: list) -> str:
        if not words:
            return ""
        lines = []
        current_line = []
        last_y = -1
        
        for w in words:
            y = w["bbox"][1]
            h = w["bbox"][3]
            if last_y == -1:
                current_line.append(w["text"])
                last_y = y
            elif abs(y - last_y) < h * 0.5:
                current_line.append(w["text"])
            else:
                lines.append(" ".join(current_line))
                current_line = [w["text"]]
                last_y = y
        if current_line:
            lines.append(" ".join(current_line))
        return "\n".join(lines)

    def _analyze_pdf(self, file_path: str, results: dict) -> dict:
        import pypdf
        try:
            reader = pypdf.PdfReader(file_path)
            pages_list = []
            
            cumulative_height = 0.0
            global_words = []
            global_raw_text_parts = []
            
            for idx, page in enumerate(reader.pages):
                page_number = idx + 1
                media_box = page.mediabox
                page_height = float(media_box.height) if media_box else 792.0
                page_width = float(media_box.width) if media_box else 612.0
                
                page_words = []
                
                def visitor_text(text, cm, tm, fontDict, fontSize):
                    if not text or not text.strip():
                        return
                    x = tm[4]
                    y = tm[5]
                    y_top = page_height - y - fontSize
                    width = len(text) * fontSize * 0.5
                    height = fontSize
                    
                    clean_text = text.strip()
                    if clean_text:
                        sub_words = clean_text.split()
                        if len(sub_words) > 1:
                            char_width = width / len(clean_text)
                            current_x = x
                            for sw in sub_words:
                                sw_len = len(sw)
                                sw_width = sw_len * char_width
                                page_words.append({
                                    "text": sw,
                                    "bbox": [round(current_x, 1), round(y_top, 1), round(sw_width, 1), round(height, 1)],
                                    "confidence": 1.0
                                })
                                current_x += (sw_len + 1) * char_width
                        else:
                            page_words.append({
                                "text": clean_text,
                                "bbox": [round(x, 1), round(y_top, 1), round(width, 1), round(height, 1)],
                                "confidence": 1.0
                            })
                
                page.extract_text(visitor_text=visitor_text)
                
                # Scanned PDF fallback: if pypdf extracts very little text, try embedded images or pdf2image
                if len(page_words) < 5:
                    page_words = []
                    # Try embedded page images first
                    if page.images:
                        for img_idx, img_obj in enumerate(page.images):
                            try:
                                from PIL import Image
                                import io
                                pil_img = Image.open(io.BytesIO(img_obj.data))
                                
                                if HAS_TESSERACT:
                                    if HAS_OPENCV:
                                        ocr_words = self._run_tesseract_with_enhancements(pil_img, pytesseract.pytesseract.tesseract_cmd)
                                    else:
                                        data_str = pytesseract.image_to_data(pil_img, output_type=pytesseract.Output.DICT)
                                        ocr_words = self._parse_tesseract_dict(data_str)
                                elif HAS_EASYOCR:
                                    # Fallback simple easyocr
                                    ocr_words = []
                                else:
                                    ocr_words = []
                                page_words.extend(ocr_words)
                            except Exception as ocr_err:
                                results["errors"].append({
                                    "plugin": self.get_name(),
                                    "severity": "warning",
                                    "message": f"Failed image OCR extraction on page {page_number}: {str(ocr_err)}"
                                })
                                
                if len(page_words) < 5:
                    # Try pdf2image rendering
                    try:
                        from pdf2image import convert_from_path
                        images = convert_from_path(file_path, first_page=page_number, last_page=page_number)
                        if images:
                            pil_img = images[0]
                            if HAS_TESSERACT:
                                if HAS_OPENCV:
                                    ocr_words = self._run_tesseract_with_enhancements(pil_img, pytesseract.pytesseract.tesseract_cmd)
                                else:
                                    data_str = pytesseract.image_to_data(pil_img, output_type=pytesseract.Output.DICT)
                                    ocr_words = self._parse_tesseract_dict(data_str)
                            else:
                                ocr_words = []
                            page_words.extend(ocr_words)
                    except Exception:
                        pass
                
                page_words.sort(key=lambda w: (w["bbox"][1], w["bbox"][0]))
                page_raw_text = self._rebuild_text_from_words(page_words)
                
                shifted_words = []
                for w in page_words:
                    sw = dict(w)
                    sw["bbox"] = [w["bbox"][0], w["bbox"][1] + cumulative_height, w["bbox"][2], w["bbox"][3]]
                    shifted_words.append(sw)
                    
                pages_list.append({
                    "page_number": page_number,
                    "raw_text": page_raw_text,
                    "words": page_words,
                    "width": page_width,
                    "height": page_height
                })
                
                global_words.extend(shifted_words)
                if page_raw_text:
                    global_raw_text_parts.append(page_raw_text)
                    
                cumulative_height += page_height
                
            results["facts"]["raw_text"] = "\n\n".join(global_raw_text_parts)
            results["facts"]["words"] = global_words
            results["facts"]["pages"] = pages_list
            
            results["indicators"].append({
                "type": "pdf_extracted",
                "description": f"Extracted native text / images from PDF file with {len(reader.pages)} pages.",
                "severity": "low"
            })
            return results
        except Exception as e:
            results["errors"].append({
                "plugin": self.get_name(),
                "severity": "error",
                "message": f"Failed to parse PDF file: {str(e)}"
            })
            return results

    def _analyze_docx(self, file_path: str, results: dict) -> dict:
        import docx
        try:
            doc = docx.Document(file_path)
            raw_text_parts = []
            for p in doc.paragraphs:
                text = p.text.strip()
                if text:
                    raw_text_parts.append(text)
            for table in doc.tables:
                table_lines = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    table_lines.append("\t".join(row_cells))
                if table_lines:
                    raw_text_parts.append("\n".join(table_lines))
                    
            raw_text = "\n\n".join(raw_text_parts)
            words = self._mock_words_from_text(raw_text)
            
            results["facts"]["raw_text"] = raw_text
            results["facts"]["words"] = words
            
            results["indicators"].append({
                "type": "docx_extracted",
                "description": "Extracted paragraphs and tables from DOCX file.",
                "severity": "low"
            })
            return results
        except Exception as e:
            results["errors"].append({
                "plugin": self.get_name(),
                "severity": "error",
                "message": f"Failed to parse DOCX document: {str(e)}"
            })
            return results

    def _analyze_pptx(self, file_path: str, results: dict) -> dict:
        import pptx
        try:
            prs = pptx.Presentation(file_path)
            pages_list = []
            cumulative_height = 0.0
            global_words = []
            global_raw_text_parts = []
            
            slide_width = prs.slide_width.inches * 96 if hasattr(prs, "slide_width") else 960.0
            slide_height = prs.slide_height.inches * 96 if hasattr(prs, "slide_height") else 720.0
            
            for idx, slide in enumerate(prs.slides):
                page_number = idx + 1
                page_words = []
                
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        x = shape.left.inches * 96 if shape.left else 20.0
                        y = shape.top.inches * 96 if shape.top else 20.0
                        w_box = shape.width.inches * 96 if shape.width else 200.0
                        
                        y_offset = y
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if not text:
                                continue
                            font_size = paragraph.font.size.pt if paragraph.font.size else 14.0
                            px_height = font_size * 1.33
                            
                            sub_words = text.split()
                            if sub_words:
                                char_width = w_box / max(len(text), 1)
                                current_x = x
                                for sw in sub_words:
                                    sw_len = len(sw)
                                    sw_width = sw_len * char_width
                                    page_words.append({
                                        "text": sw,
                                        "bbox": [round(current_x, 1), round(y_offset, 1), round(sw_width, 1), round(px_height, 1)],
                                        "confidence": 1.0
                                    })
                                    current_x += (sw_len + 1) * char_width
                            y_offset += px_height + 5
                            
                page_words.sort(key=lambda w: (w["bbox"][1], w["bbox"][0]))
                page_raw_text = self._rebuild_text_from_words(page_words)
                
                shifted_words = []
                for w in page_words:
                    sw = dict(w)
                    sw["bbox"] = [w["bbox"][0], w["bbox"][1] + cumulative_height, w["bbox"][2], w["bbox"][3]]
                    shifted_words.append(sw)
                    
                pages_list.append({
                    "page_number": page_number,
                    "raw_text": page_raw_text,
                    "words": page_words,
                    "width": slide_width,
                    "height": slide_height
                })
                global_words.extend(shifted_words)
                if page_raw_text:
                    global_raw_text_parts.append(page_raw_text)
                cumulative_height += slide_height
                
            results["facts"]["raw_text"] = "\n\n".join(global_raw_text_parts)
            results["facts"]["words"] = global_words
            results["facts"]["pages"] = pages_list
            
            results["indicators"].append({
                "type": "pptx_extracted",
                "description": f"Extracted native text from PPTX Presentation with {len(prs.slides)} slides.",
                "severity": "low"
            })
            return results
        except Exception as e:
            results["errors"].append({
                "plugin": self.get_name(),
                "severity": "error",
                "message": f"Failed to parse PPTX presentation: {str(e)}"
            })
            return results

    def _analyze_xlsx(self, file_path: str, results: dict) -> dict:
        import pandas as pd
        try:
            xl = pd.ExcelFile(file_path)
            pages_list = []
            cumulative_height = 0.0
            global_words = []
            global_raw_text_parts = []
            
            sheet_width = 800.0
            sheet_height = 600.0
            
            for idx, sheet_name in enumerate(xl.sheet_names):
                page_number = idx + 1
                df = xl.parse(sheet_name)
                df_str = df.to_string(index=False)
                
                page_words = []
                lines = df_str.split("\n")
                
                y_offset = 20.0
                for line in lines:
                    line_words = line.strip().split()
                    if not line_words:
                        y_offset += 25.0
                        continue
                    x_offset = 20.0
                    for lw in line_words:
                        w_len = len(lw) * 8.0
                        page_words.append({
                            "text": lw,
                            "bbox": [round(x_offset, 1), round(y_offset, 1), round(w_len, 1), 20.0],
                            "confidence": 1.0
                        })
                        x_offset += w_len + 12.0
                    y_offset += 30.0
                    
                page_raw_text = df_str
                shifted_words = []
                for w in page_words:
                    sw = dict(w)
                    sw["bbox"] = [w["bbox"][0], w["bbox"][1] + cumulative_height, w["bbox"][2], w["bbox"][3]]
                    shifted_words.append(sw)
                    
                pages_list.append({
                    "page_number": page_number,
                    "page_name": sheet_name,
                    "raw_text": page_raw_text,
                    "words": page_words,
                    "width": sheet_width,
                    "height": sheet_height
                })
                global_words.extend(shifted_words)
                if page_raw_text:
                    global_raw_text_parts.append(page_raw_text)
                cumulative_height += sheet_height
                
            results["facts"]["raw_text"] = "\n\n".join(global_raw_text_parts)
            results["facts"]["words"] = global_words
            results["facts"]["pages"] = pages_list
            
            results["indicators"].append({
                "type": "xlsx_extracted",
                "description": f"Extracted sheets data from Excel file with {len(xl.sheet_names)} sheets.",
                "severity": "low"
            })
            try:
                xl.close()
            except Exception:
                pass
            return results
        except Exception as e:
            results["errors"].append({
                "plugin": self.get_name(),
                "severity": "error",
                "message": f"Failed to parse Excel file: {str(e)}"
            })
            return results

    def _analyze_text_file(self, file_path: str, results: dict) -> dict:
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read().strip()
                
            words = self._mock_words_from_text(content)
            results["facts"]["raw_text"] = content
            results["facts"]["words"] = words
            
            results["indicators"].append({
                "type": "text_extracted",
                "description": f"Loaded plain text file: {os.path.basename(file_path)}",
                "severity": "low"
            })
            return results
        except Exception as e:
            results["errors"].append({
                "plugin": self.get_name(),
                "severity": "error",
                "message": f"Failed to read text file: {str(e)}"
            })
            return results
