import os
import tempfile
import unittest
import shutil
from PIL import Image
from image_extractor.extractor import ImageInfoExtractor


class TestImageTextLayoutExtractor(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        # Create a temporary directory for test assets
        cls.test_dir = tempfile.mkdtemp()
        
        # 1. Simple dummy image
        cls.dummy_path = os.path.join(cls.test_dir, "dummy.png")
        Image.new("RGB", (100, 100), color="white").save(cls.dummy_path)

        # 2. Document page scan simulation (PNG + Sidecar OCR Text file)
        cls.doc_path = os.path.join(cls.test_dir, "test_doc.png")
        Image.new("RGB", (120, 150), color="white").save(cls.doc_path)
        
        # Write sidecar .txt file next to it
        cls.doc_txt_path = os.path.join(cls.test_dir, "test_doc.txt")
        cls.doc_text = (
            "Mrs. Russell shook her head and cried.\n"
            "\"Never mind, darling, you shall have them one day.\"\n\n"
            "\"I know it sounds lovely,\" said Bob. \"Southend is nice, but Uncle Edward is waiting.\"\n"
            "Here is an IP address: 192.168.1.100 and a URL: https://example.com/api\n"
            "when in rome, do as the romans do\n"
            "the customer is always right\n"
            "east, west, home's best"
        )
        with open(cls.doc_txt_path, "w", encoding="utf-8") as f:
            f.write(cls.doc_text)

        # 3. Create PDF
        cls.pdf_path = os.path.join(cls.test_dir, "test.pdf")
        from reportlab.pdfgen import canvas
        c = canvas.Canvas(cls.pdf_path)
        c.drawString(100, 750, "Hello World from reportlab PDF")
        c.drawString(100, 700, "This is the second line of text.")
        c.save()

        # 4. Create DOCX
        cls.docx_path = os.path.join(cls.test_dir, "test.docx")
        import docx
        doc = docx.Document()
        doc.add_paragraph("Hello from python-docx Word document")
        doc.add_paragraph("Second paragraph in the document.")
        doc.save(cls.docx_path)

        # 5. Create PPTX
        cls.pptx_path = os.path.join(cls.test_dir, "test.pptx")
        import pptx
        prs = pptx.Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Hello from python-pptx PowerPoint"
        txBox = slide.shapes.add_textbox(100, 100, 200, 50)
        tf = txBox.text_frame
        tf.text = "This is a slide body paragraph."
        prs.save(cls.pptx_path)

        # 6. Create XLSX
        cls.xlsx_path = os.path.join(cls.test_dir, "test.xlsx")
        import pandas as pd
        df = pd.DataFrame({
            "Name": ["Alice", "Bob"],
            "Age": [25, 30]
        })
        df.to_excel(cls.xlsx_path, index=False)

        # 7. Create TXT
        cls.txt_path = os.path.join(cls.test_dir, "test.txt")
        with open(cls.txt_path, "w", encoding="utf-8") as f:
            f.write("Hello from a simple plain text file.")

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls.test_dir)

    def test_invalid_paths(self):
        with self.assertRaises(FileNotFoundError):
            ImageInfoExtractor("missing_file.jpg")
        with self.assertRaises(ValueError):
            ImageInfoExtractor(self.test_dir)

    def test_text_and_layout_extraction(self):
        extractor = ImageInfoExtractor(self.doc_path)
        results = extractor.extract_all()
        
        # Verify sidecar loaded indicator is logged
        self.assertTrue(any(ind["type"] == "ocr_sidecar_loaded" for ind in results["indicators"]))
        
        # Verify paragraph and layout reconstruction
        facts = results["facts"]
        self.assertIn("paragraphs", facts)
        self.assertGreaterEqual(len(facts["paragraphs"]), 2)
        
        # Ensure we have columns information
        self.assertEqual(facts["paragraphs"][0]["column"], 1)
        
        stats = facts["statistics"]
        self.assertGreater(stats["word_count"], 15)
        self.assertGreater(stats["line_count"], 3)
        
        # Verify language detection and classification
        assessments = results["assessments"]
        self.assertIn("language_detection", assessments)
        self.assertIn("English", assessments["language_detection"]["language"])
        
        self.assertIn("document_classification", assessments)
        # Should detect Scanned Printed Book Page due to narrative keywords
        self.assertEqual(assessments["document_classification"]["document_type"], "Scanned Printed Book Page")
        self.assertEqual(assessments["document_classification"]["content_type"], "Narrative Text / Prose")

    def test_pdf_extraction(self):
        extractor = ImageInfoExtractor(self.pdf_path)
        results = extractor.extract_all()
        facts = results["facts"]
        self.assertIn("pages", facts)
        self.assertEqual(len(facts["pages"]), 1)
        self.assertIn("reportlab", facts["raw_text"])
        self.assertEqual(facts["statistics"]["word_count"], 12)

    def test_docx_extraction(self):
        extractor = ImageInfoExtractor(self.docx_path)
        results = extractor.extract_all()
        facts = results["facts"]
        self.assertIn("docx", facts["raw_text"])
        self.assertEqual(facts["statistics"]["word_count"], 10)

    def test_pptx_extraction(self):
        extractor = ImageInfoExtractor(self.pptx_path)
        results = extractor.extract_all()
        facts = results["facts"]
        self.assertIn("pages", facts)
        self.assertEqual(len(facts["pages"]), 1)
        self.assertIn("pptx", facts["raw_text"])
        self.assertEqual(facts["statistics"]["word_count"], 10)

    def test_xlsx_extraction(self):
        extractor = ImageInfoExtractor(self.xlsx_path)
        results = extractor.extract_all()
        facts = results["facts"]
        self.assertIn("pages", facts)
        self.assertEqual(len(facts["pages"]), 1)
        self.assertIn("Alice", facts["raw_text"])
        self.assertEqual(facts["statistics"]["word_count"], 6)

    def test_txt_extraction(self):
        extractor = ImageInfoExtractor(self.txt_path)
        results = extractor.extract_all()
        facts = results["facts"]
        self.assertIn("simple", facts["raw_text"])
        self.assertEqual(facts["statistics"]["word_count"], 7)


if __name__ == "__main__":
    unittest.main()
